"""
Generate an HTML visual map of all 7 slides showing exact shape positions.
This lets us visually verify nothing overlaps in the browser.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

ppt_path = r'c:\Users\sidda\Downloads\TBIE_CODE\TBIE_Solution_Presentation.pptx'
prs = Presentation(ppt_path)

SLIDE_W = prs.slide_width / 914400
SLIDE_H = prs.slide_height / 914400
SCALE = 130  # pixels per inch

html = """<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>PPT Layout Visual Audit</title>
<style>
  body { background: #1a1a2e; color: #fff; font-family: 'Segoe UI', sans-serif; padding: 20px; }
  h1 { color: #e66c37; text-align: center; }
  .slide-container { margin: 30px auto; position: relative; }
  .slide-label { font-size: 18px; font-weight: bold; color: #e66c37; margin-bottom: 8px; }
  .slide { position: relative; border: 2px solid #444; overflow: hidden; margin-bottom: 40px; }
  .shape { position: absolute; border: 1px solid rgba(255,255,255,0.15); overflow: hidden;
           font-size: 7px; line-height: 1.2; color: #ccc; padding: 1px 2px; box-sizing: border-box; }
  .shape.text { background: rgba(46, 125, 50, 0.15); border-color: rgba(46, 125, 50, 0.4); }
  .shape.table { background: rgba(230, 108, 55, 0.2); border-color: rgba(230, 108, 55, 0.5); }
  .shape.rect { background: rgba(100, 100, 100, 0.08); border-color: rgba(100,100,100,0.2); }
  .shape:hover { z-index: 999; border-color: #fff !important; background: rgba(255,255,255,0.15) !important; }
  .shape:hover::after { content: attr(data-info); position: absolute; bottom: -20px; left: 0;
    background: #000; color: #fff; padding: 2px 5px; font-size: 9px; white-space: nowrap; z-index: 1000; }
</style></head><body>
<h1>PPT Layout Visual Audit — All 7 Slides</h1>
<p style="text-align:center;color:#888;">Green = text | Orange = table | Gray = empty rect/border. Hover for details.</p>
"""

for si, slide in enumerate(prs.slides):
    # Determine background color
    bg = "#fdfbf7"  # light
    if si == 0 or si == 5:
        bg = "#1b212c"  # dark slides
    
    pw = SLIDE_W * SCALE
    ph = SLIDE_H * SCALE
    
    html += f'<div class="slide-container" style="width:{pw}px;">\n'
    html += f'<div class="slide-label">Slide {si+1}</div>\n'
    html += f'<div class="slide" style="width:{pw}px; height:{ph}px; background:{bg};">\n'
    
    for idx, shape in enumerate(slide.shapes):
        left = (shape.left / 914400 if shape.left else 0) * SCALE
        top = (shape.top / 914400 if shape.top else 0) * SCALE
        w = (shape.width / 914400 if shape.width else 0) * SCALE
        h = (shape.height / 914400 if shape.height else 0) * SCALE
        
        text = ""
        cls = "rect"
        if shape.has_table:
            cls = "table"
            t = shape.table
            text = f"TABLE[{len(t.rows)}x{len(t.columns)}]"
        elif shape.has_text_frame:
            full = shape.text_frame.text.strip()
            if full:
                cls = "text"
                text = full[:80].replace('<','&lt;').replace('>','&gt;').replace('"','&quot;')
        
        info = f"#{idx} ({left/SCALE:.2f},{top/SCALE:.2f})-({(left+w)/SCALE:.2f},{(top+h)/SCALE:.2f})"
        
        html += f'  <div class="shape {cls}" style="left:{left:.0f}px;top:{top:.0f}px;width:{w:.0f}px;height:{h:.0f}px;" '
        html += f'data-info="{info}" title="{text[:60]}">{text[:40]}</div>\n'
    
    html += '</div></div>\n'

html += '</body></html>'

out_path = r'c:\Users\sidda\Downloads\TBIE_CODE\TBIE_CODE\layout_visual.html'
with open(out_path, 'w', encoding='utf-8') as f:
    f.write(html)
print(f"Saved: {out_path}")
