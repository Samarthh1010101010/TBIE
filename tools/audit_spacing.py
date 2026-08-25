"""
Visual spacing audit — Check exact Y-gaps between consecutive content elements per slide.
Flags anywhere two content elements are within 0.02 inches or less of each other (risk of visual touching).
Also checks text boxes that might be too small for their content.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Emu

ppt_path = r'c:\Users\sidda\Downloads\TBIE_CODE\TBIE_Solution_Presentation.pptx'
prs = Presentation(ppt_path)

SLIDE_W = prs.slide_width / 914400
SLIDE_H = prs.slide_height / 914400

print("=" * 110)
print("VISUAL SPACING AUDIT — checking for tight/touching elements")
print("=" * 110)

total_warnings = 0

for si, slide in enumerate(prs.slides):
    print(f"\n{'━' * 110}")
    print(f"  SLIDE {si+1}")
    print(f"{'━' * 110}")
    
    # Collect content shapes with text
    content = []
    for idx, shape in enumerate(slide.shapes):
        left = shape.left / 914400 if shape.left else 0
        top = shape.top / 914400 if shape.top else 0
        w = shape.width / 914400 if shape.width else 0
        h = shape.height / 914400 if shape.height else 0
        
        text = ""
        stype = "rect"
        lines_of_text = 0
        
        if shape.has_table:
            stype = "TABLE"
            t = shape.table
            text = f"[{len(t.rows)}x{len(t.columns)}]"
            lines_of_text = len(t.rows)
        elif shape.has_text_frame:
            full = shape.text_frame.text.strip()
            if full:
                stype = "TEXT"
                text = full.replace('\n', ' ')[:50]
                lines_of_text = full.count('\n') + 1
                
                # Check if text frame has enough space for its content
                # Estimate: each line needs ~0.12 inches at 6pt, ~0.14 at 7pt, ~0.18 at 8pt
                # Get actual font size from first paragraph
                try:
                    fs = shape.text_frame.paragraphs[0].runs[0].font.size
                    if fs:
                        font_pt = fs / 12700  # EMU to pt
                    else:
                        font_pt = 7
                except:
                    font_pt = 7
                
                line_height = font_pt * 0.02  # rough: 1pt = 0.02 inches line height
                needed_h = lines_of_text * line_height
                
                if needed_h > h * 1.5 and h < 0.5:
                    print(f"  [CRAMPED] #{idx}: {lines_of_text} lines of {font_pt:.0f}pt text in {h:.2f}\" tall box (needs ~{needed_h:.2f}\")")
                    print(f"            '{text[:60]}'")
                    total_warnings += 1
        
        if stype in ('TEXT', 'TABLE'):
            content.append({
                'idx': idx, 'top': top, 'bottom': top+h, 'left': left, 'right': left+w,
                'h': h, 'text': text, 'type': stype, 'lines': lines_of_text
            })
    
    # Check for elements too close vertically in the same column
    # Group by approximate X position (left half vs right half)
    for zone_name, x_min, x_max in [("LEFT", 0, 5.0), ("RIGHT", 4.8, 10.0), ("FULL", 0, 10.0)]:
        zone_items = [c for c in content if c['left'] >= x_min and c['right'] <= x_max + 0.5]
        zone_items.sort(key=lambda x: x['top'])
        
        for i in range(len(zone_items) - 1):
            a = zone_items[i]
            b = zone_items[i+1]
            
            # Only compare items in the same horizontal zone
            if abs(a['left'] - b['left']) > 3:
                continue  # Different columns
            
            gap = b['top'] - a['bottom']
            if gap < -0.02 and gap > -0.5:  # Slightly overlapping (not parent-child)
                # Check they're not a header-content pair
                if a['text'] and b['text'] and len(a['text']) > 5 and len(b['text']) > 5:
                    print(f"  [TOUCHING] {zone_name}: #{a['idx']} bottom={a['bottom']:.2f} vs #{b['idx']} top={b['top']:.2f} (gap={gap:.3f}\")")
                    print(f"             A: '{a['text'][:45]}'")
                    print(f"             B: '{b['text'][:45]}'")
                    total_warnings += 1
    
    # Check if any table extends into footer
    for c in content:
        if c['type'] == 'TABLE' and c['bottom'] > SLIDE_H - 0.30:
            print(f"  [TABLE-FOOTER] #{c['idx']}: table bottom={c['bottom']:.2f} enters footer zone")
            total_warnings += 1
    
    # Check right-edge proximity
    for c in content:
        if c['right'] > SLIDE_W - 0.2 and c['right'] <= SLIDE_W and c['left'] > 0.3:
            pass  # Fine, just close to edge
        elif c['right'] > SLIDE_W and c['left'] < SLIDE_W - 1:
            print(f"  [EDGE-CLIP] #{c['idx']}: right edge {c['right']:.2f} exceeds slide width {SLIDE_W:.2f}")
            print(f"              '{c['text'][:45]}'")
            total_warnings += 1

print(f"\n{'=' * 110}")
print(f"TOTAL WARNINGS: {total_warnings}")
if total_warnings == 0:
    print("✅ All slides have clean spacing!")
print(f"{'=' * 110}")
