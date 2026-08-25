"""
PPT Layout Audit v2 — Only flag REAL overlaps between content-bearing shapes.
Excludes: 
  - Background rect + text-on-top (intended parent-child layering)
  - Footer components
  - Decorative circles (Slide 1)
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

ppt_path = r'c:\Users\sidda\Downloads\TBIE_CODE\TBIE_Solution_Presentation.pptx'
prs = Presentation(ppt_path)

SLIDE_W = prs.slide_width / 914400
SLIDE_H = prs.slide_height / 914400

print(f"Slide: {SLIDE_W:.2f} x {SLIDE_H:.2f} inches")
print("=" * 110)

total_real_issues = 0

for si, slide in enumerate(prs.slides):
    print(f"\n{'━' * 110}")
    print(f"  SLIDE {si+1}")
    print(f"{'━' * 110}")
    
    shapes_info = []
    for idx, shape in enumerate(slide.shapes):
        left = shape.left / 914400 if shape.left is not None else 0
        top = shape.top / 914400 if shape.top is not None else 0
        w = shape.width / 914400 if shape.width is not None else 0
        h = shape.height / 914400 if shape.height is not None else 0
        
        text = ""
        has_content = False
        shape_type = "rect"  # default = empty shape/rectangle
        
        if shape.has_table:
            shape_type = "TABLE"
            has_content = True
            t = shape.table
            text = f"[{len(t.rows)}x{len(t.columns)}] " + t.cell(0,0).text[:30]
        elif shape.has_text_frame:
            full = shape.text_frame.text.strip()
            if full:
                shape_type = "TEXT"
                has_content = True
                text = full.replace('\n', ' ')[:60]
            else:
                shape_type = "rect"  # empty text frame = border/background
        
        shapes_info.append({
            'idx': idx, 'left': left, 'top': top, 'w': w, 'h': h,
            'right': left+w, 'bottom': top+h,
            'text': text, 'type': shape_type, 'has_content': has_content
        })
    
    issues = []
    
    # ── Only check content vs content overlaps ──
    content_shapes = [s for s in shapes_info if s['has_content']]
    
    for i in range(len(content_shapes)):
        a = content_shapes[i]
        for j in range(i+1, len(content_shapes)):
            b = content_shapes[j]
            
            # Skip footer elements overlapping each other (they're in different positions)
            if a['top'] > SLIDE_H - 0.35 and b['top'] > SLIDE_H - 0.35:
                continue
            
            # Skip decorative (slide 1 huge shapes)
            if a['w'] > 5 or b['w'] > 5:
                continue
            
            # Check overlap
            if (a['left'] < b['right'] and a['right'] > b['left'] and
                a['top'] < b['bottom'] and a['bottom'] > b['top']):
                
                ox = min(a['right'], b['right']) - max(a['left'], b['left'])
                oy = min(a['bottom'], b['bottom']) - max(a['top'], b['top'])
                overlap_area = ox * oy
                
                a_area = max(a['w'] * a['h'], 0.01)
                b_area = max(b['w'] * b['h'], 0.01)
                pct = max(overlap_area/a_area, overlap_area/b_area) * 100
                
                # Only flag if the overlap is significant (>3% of either)
                if pct > 3:
                    # Check if one is clearly a container for the other (parent-child)
                    # Parent = no content or section header, child = detailed text
                    a_is_parent = (a['type'] == 'TEXT' and len(a['text']) < 30 and 
                                   a['left'] <= b['left'] and a['top'] <= b['top'] and
                                   a['right'] >= b['right'] and a['bottom'] >= b['bottom'])
                    b_is_parent = (b['type'] == 'TEXT' and len(b['text']) < 30 and 
                                   b['left'] <= a['left'] and b['top'] <= a['top'] and
                                   b['right'] >= a['right'] and b['bottom'] >= a['bottom'])
                    
                    if a_is_parent or b_is_parent:
                        continue  # Intentional header-over-box
                    
                    # Check if this is a section header just above content (tiny Y overlap)
                    if oy < 0.06:
                        continue  # Negligible border overlap
                    
                    severity = "OVERLAP" if pct < 30 else "MAJOR"
                    issues.append(
                        f"  [{severity}] #{a['idx']} vs #{b['idx']}: "
                        f"overlap={overlap_area:.2f}in² ({pct:.0f}%)\n"
                        f"      A: Y={a['top']:.2f}-{a['bottom']:.2f} X={a['left']:.2f}-{a['right']:.2f} '{a['text'][:45]}'\n"
                        f"      B: Y={b['top']:.2f}-{b['bottom']:.2f} X={b['left']:.2f}-{b['right']:.2f} '{b['text'][:45]}'"
                    )
    
    # ── Check content going below footer zone ──
    for s in content_shapes:
        if s['top'] < SLIDE_H - 0.35 and s['bottom'] > SLIDE_H - 0.25:
            issues.append(
                f"  [CLIP-BOTTOM] #{s['idx']}: bottom={s['bottom']:.2f} enters footer zone (>{SLIDE_H-0.28:.2f})\n"
                f"      '{s['text'][:50]}'"
            )
    
    # ── Check content going off right edge ──
    for s in content_shapes:
        if s['w'] < 5 and s['right'] > SLIDE_W + 0.02:
            issues.append(
                f"  [OFF-RIGHT] #{s['idx']}: right={s['right']:.2f} > slide width {SLIDE_W:.2f}\n"
                f"      '{s['text'][:50]}'"
            )
    
    # ── Print content layout map ──
    print(f"  {'#':>3} {'Type':>5} {'Y-top':>5} {'Y-bot':>5} {'X-lft':>5} {'X-rgt':>5}  Content")
    for s in content_shapes:
        print(f"  {s['idx']:>3} {s['type']:>5} {s['top']:>5.2f} {s['bottom']:>5.2f} {s['left']:>5.2f} {s['right']:>5.2f}  {s['text'][:65]}")
    
    if issues:
        print(f"\n  ⚠️  {len(issues)} REAL ISSUE(S):")
        for iss in issues:
            print(iss)
        total_real_issues += len(issues)
    else:
        print(f"\n  ✅ Clean layout")

print(f"\n{'=' * 110}")
print(f"REAL ISSUES TOTAL: {total_real_issues}")
print(f"{'=' * 110}")
