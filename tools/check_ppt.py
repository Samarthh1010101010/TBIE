import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_ppt(file_path):
    try:
        prs = Presentation(file_path)
    except Exception as e:
        print(f"Error loading {file_path}: {e}")
        return

    print(f"--- Analysis of {file_path} ---")
    print(f"Total Slides: {len(prs.slides)}")
    print(f"Dimensions (Inches): {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f}")
    
    total_shapes = 0
    text_boxes = 0
    tables = 0
    
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}:")
        slide_shapes = len(slide.shapes)
        slide_tables = 0
        total_shapes += slide_shapes
        
        for shape in slide.shapes:
            if shape.has_table:
                slide_tables += 1
                tables += 1
            elif shape.has_text_frame:
                text_boxes += 1
                
        print(f"  Shapes: {slide_shapes}, Tables: {slide_tables}")
        
        # Check text length / overflow heuristically
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.replace('\n', ' ').strip()
                if text:
                    y_pos = shape.top / 914400
                    h = shape.height / 914400
                    if y_pos + h > 5.6:
                        print(f"  [!] Potential Overflow near bottom: '{text[:40]}...' (Y: {y_pos:.2f}, H: {h:.2f})")
    
    print(f"\nTotal Shapes: {total_shapes}, Total Text Boxes: {text_boxes}, Total Tables: {tables}\n")

analyze_ppt(r'c:\Users\sidda\Downloads\TBIE_CODE\solution_presentation_template.pptx')
analyze_ppt(r'c:\Users\sidda\Downloads\TBIE_CODE\TBIE_Solution_Presentation.pptx')
