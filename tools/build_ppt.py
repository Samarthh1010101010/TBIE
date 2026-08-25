"""
TBIE Presentation Builder — Full 7-Slide Deck (Updated with verified content)
Total slides: 7
"""

import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette (UNCHANGED) ──
LIGHT_BG      = RGBColor(253, 251, 247)  # #FDFBF7 Off-white
DARK_TEXT     = RGBColor(27, 33, 44)     # #1B212C Slate navy
BODY_TEXT     = RGBColor(89, 89, 89)     # #595959 Dark gray
ACCENT_ORANGE = RGBColor(230, 108, 55)   # #E66C37 Bright orange
LIGHT_ORANGE  = RGBColor(255, 204, 188)  # #FFCCBC For box outlines
FOOTER_BG     = RGBColor(38, 34, 32)     # #262220 Very dark brown
WHITE         = RGBColor(255, 255, 255)
LIGHT_GRAY    = RGBColor(200, 200, 200)

# Semantic colors for text
GREEN_OK      = RGBColor(46, 125, 50)
RED_ALERT     = RGBColor(211, 47, 47)
ACCENT_TEAL   = RGBColor(38, 166, 154)

SH = 5.625  # Slide height in inches

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid(); fill.fore_color.rgb = color

def add_tb(slide, left, top, width, height, text, sz=14, bold=False,
           color=BODY_TEXT, align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font; p.alignment = align
    return tb

def add_ml(slide, left, top, width, height, lines, default_sz=10,
           default_color=BODY_TEXT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, col, sz = item, False, default_color, default_sz
        else:
            txt = item[0]; bld = item[1] if len(item)>1 else False
            col = item[2] if len(item)>2 else default_color
            sz = item[3] if len(item)>3 else default_sz
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = txt
        r.font.size = Pt(sz); r.font.bold = bld
        r.font.color.rgb = col; r.font.name = font
    return tb

def add_rect(slide, left, top, width, height, fill=None, line=LIGHT_ORANGE, thick_top=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()

    if line:
        s.line.color.rgb = line; s.line.width = Pt(1)
    else:
        s.line.fill.background()

    if thick_top:
        tb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(0.08))
        tb.fill.solid(); tb.fill.fore_color.rgb = ACCENT_ORANGE
        tb.line.fill.background()
    return s

def add_footer(slide, page_num, total_pages=7):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(SH - 0.28), Inches(10), Inches(0.28))
    s.fill.solid(); s.fill.fore_color.rgb = FOOTER_BG; s.line.fill.background()
    add_tb(slide, 0.35, SH-0.27, 3.5, 0.25, 'KOBIE LAUNCHPAD AI HACKATHON',
           sz=7, bold=True, color=ACCENT_ORANGE)
    add_tb(slide, 8.8, SH-0.27, 0.9, 0.25, f'{page_num} / {total_pages}',
           sz=7, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

def add_table(slide, left, top, width, height, data, col_w=None, font_size=6, theme='light'):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    t = ts.table
    if col_w:
        for i, w in enumerate(col_w):
            t.columns[i].width = Inches(w)
    for r in range(rows):
        for c in range(cols):
            cell = t.cell(r, c); cell.text = str(data[r][c])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size); p.font.name = 'Calibri'
                p.font.bold = (r == 0)
                if theme == 'light':
                    p.font.color.rgb = WHITE if r==0 else DARK_TEXT
                else:
                    p.font.color.rgb = WHITE

            cell.fill.solid()
            if theme == 'light':
                if r == 0:
                    cell.fill.fore_color.rgb = ACCENT_ORANGE
                elif r % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(245, 245, 245)
                else:
                    cell.fill.fore_color.rgb = WHITE
            else:
                if r == 0:
                    cell.fill.fore_color.rgb = ACCENT_ORANGE
                elif r % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(45, 55, 70)
                else:
                    cell.fill.fore_color.rgb = RGBColor(35, 45, 60)
    return ts

# ══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(10)
prs.slide_height = Inches(SH)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 — Title (Dark Theme preserved)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
TITLE_BG = RGBColor(0x1B, 0x21, 0x2C)
set_slide_bg(sl, TITLE_BG)

circle1 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.0), Inches(-3.0), Inches(7.5), Inches(7.5))
circle1.fill.solid(); circle1.fill.fore_color.rgb = RGBColor(0xC8, 0x69, 0x44); circle1.line.fill.background()

circle2 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), Inches(2.5), Inches(6.0), Inches(6.0))
circle2.fill.solid(); circle2.fill.fore_color.rgb = RGBColor(0x99, 0x4B, 0x32); circle2.line.fill.background()

BOX_BG = RGBColor(0x4B, 0x3D, 0x3A)

add_tb(sl, 0.5, 0.55, 6, 0.3, 'KOBIE LAUNCHPAD AI HACKATHON', sz=11, bold=True, color=ACCENT_ORANGE)
add_tb(sl, 0.5, 0.85, 6, 0.3, 'Solution Presentation', sz=14, color=WHITE)

add_tb(sl, 0.5, 1.4, 7, 0.6, 'TBIE', sz=44, bold=True, color=WHITE, font='Georgia')
add_tb(sl, 0.5, 2.1, 7, 0.35, 'Temporal Behavioural Intelligence Engine', sz=16, color=WHITE)
add_tb(sl, 0.5, 2.6, 7, 0.3, 'Track C: Dynamic Member Segmentation & Activation', sz=12, color=WHITE)

add_rect(sl, 0.5, 3.2, 5.0, 1.3, fill=BOX_BG, line=None)
add_ml(sl, 0.7, 3.32, 4.6, 1.1, [
    ('Team Members', True, ACCENT_ORANGE, 11),
    ('Samarth Vinod Hosalli (PES1UG23AM261)', False, WHITE, 9),
    ('Siddarth Reddy (PES1UG23AM300)', False, WHITE, 9),
])
add_tb(sl, 0.5, 4.75, 4, 0.25, 'PES University', sz=10, color=WHITE)

add_rect(sl, 6.5, 1.5, 3.0, 2.7, fill=TITLE_BG, line=ACCENT_ORANGE)
stats = [('Macro F1 (Test)','0.8138',WHITE),('Segments','5',ACCENT_ORANGE),
         ('States','10',WHITE),('Members','500,000',ACCENT_ORANGE),('Features Built','119',WHITE)]
for i,(lbl,val,col) in enumerate(stats):
    y = 1.6 + i*0.48
    add_tb(sl, 6.7, y, 2.5, 0.22, val, sz=18, bold=True, color=col)
    add_tb(sl, 6.7, y+0.22, 2.5, 0.16, lbl, sz=8, color=WHITE)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 — Problem & Our Approach (Light Theme)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, LIGHT_BG); add_footer(sl, 1)
add_tb(sl, 0.5, 0.25, 7, 0.5, 'Problem & Our Approach', sz=28, bold=True, color=DARK_TEXT, font='Georgia')

# ── LEFT: THE CORE CHALLENGE ──
add_tb(sl, 0.5, 0.85, 3, 0.25, 'THE CORE CHALLENGE', sz=11, bold=True, color=ACCENT_ORANGE)
add_rect(sl, 0.5, 1.1, 4.3, 2.4, fill=None, line=LIGHT_ORANGE)
add_ml(sl, 0.65, 1.15, 4.0, 2.3, [
    ('The real challenge is not segmentation — it is separating signal from noise where the most important behaviours are the most ambiguous.', False, DARK_TEXT, 7),
    ('', False, BODY_TEXT, 4),
    ('What makes this hard:', True, ACCENT_ORANGE, 7.5),
    ('1. Members form a continuous gradient, not discrete groups. No density gaps. HDBSCAN labelled 71–85% as noise.', False, BODY_TEXT, 6.5),
    ('2. Temporal leakage is invisible. members.parquet has 3 columns (lifetime_points_earned, lifetime_points_redeemed, current_point_balance) reflecting Dec 2025 totals, not as-of-observation. We reconstructed all point features from raw transactions.', False, BODY_TEXT, 6.5),
    ('3. Segment and state are two different problems that look like one. A Growth Builder in Lapse Risk and a Silent Accumulator in Lapse Risk both show "low recent purchases" — but need different interventions.', False, BODY_TEXT, 6.5),
    ('', False, BODY_TEXT, 4),
    ('Constraints:', True, DARK_TEXT, 7.5),
    ('• Frozen models — no refit at inference, same model scores any future date identically', False, BODY_TEXT, 6.5),
    ('• Single-command reproducibility: python pipeline.py --observation_date 2025-12-31', False, BODY_TEXT, 6.5),
    ('• No ground truth — every threshold justified by data distribution + business logic', False, BODY_TEXT, 6.5),
])

# ── RIGHT: OUR KEY INSIGHT ──
add_tb(sl, 5.2, 0.85, 3, 0.25, 'OUR KEY INSIGHT', sz=11, bold=True, color=ACCENT_ORANGE)
add_rect(sl, 5.2, 1.1, 4.3, 2.4, fill=None, line=LIGHT_ORANGE)
add_ml(sl, 5.35, 1.15, 4.0, 2.3, [
    ('A member\'s segment tells you who they are over months. Their state tells you what to do this week.', True, DARK_TEXT, 7.5),
    ('Acting on segment alone misses the moment. Acting on state alone misses the trajectory. You need both.', False, BODY_TEXT, 7),
    ('', False, BODY_TEXT, 4),
    ('How this addresses the hard parts:', True, ACCENT_ORANGE, 7.5),
    ('• Continuous gradient → K-Means as a partitioning tool (not a cluster-finder), validated by actionability, not silhouette alone.', False, BODY_TEXT, 6.5),
    ('• Leakage → Every feature reconstructed from raw transactions + events anchored to obs_date. No static member-level aggregates ever used.', False, BODY_TEXT, 6.5),
    ('• Segment-vs-state → Two-layer architecture: K-Means assigns long-run archetypes; 9 vectorised business rules assign short-run postures; XGBoost predicts 30-day segment transitions.', False, BODY_TEXT, 6.5),
    ('', False, BODY_TEXT, 4),
    ('Alternatives tried & rejected:', True, DARK_TEXT, 7.5),
    ('• HDBSCAN: 71–85% noise. No density gaps to exploit.', False, BODY_TEXT, 6.5),
    ('• GMM: Soft memberships not operationalisable ("40% Silent Acc, 35% Growth Builder").', False, BODY_TEXT, 6.5),
    ('• k=7: Two centroids within 0.4 PCA distance — functionally identical.', False, BODY_TEXT, 6.5),
    ('• Neural embeddings: No interpretable trace when state is misclassified.', False, BODY_TEXT, 6.5),
    ('• Learned state classification: States must be auditable. "recency > 60 AND app_open > 0" is explainable. A learned classifier is not.', False, BODY_TEXT, 6.5),
])

# ── BOTTOM: DATA QUALITY TABLE ──
add_rect(sl, 0.5, 3.65, 9.0, 1.55, fill=None, line=LIGHT_ORANGE, thick_top=True)
add_tb(sl, 0.65, 3.75, 8, 0.2, 'DATA QUALITY — 8 ISSUES FOUND & RESOLVED IN EDA', sz=9, bold=True, color=ACCENT_ORANGE)
dq = [
    ['#', 'Issue', 'Scale', 'Resolution'],
    ['1', 'Ghost member IDs (MBR_GHOST_#)', '88,717 orphaned IDs (17.7%)', 'Excluded via spine LEFT JOIN'],
    ['2', 'Mixed datetime formats (DD-MM / YYYY-MM)', '~7.4% of date rows', 'Dual-pass parser (ISO first, day-first fallback)'],
    ['3', 'Mixed-case transaction types', '24,738 rows', '.str.lower() normalisation'],
    ['4', 'Duplicate engagement events', '92,748 rows (46,374 groups)', 'drop_duplicates on [member_id, event_date, event_type]'],
    ['5', 'Session duration outliers (max 48h)', '18,663 rows > 4h', 'Clipped at 14,400s (4 hours)'],
    ['6', 'Leakage columns (points earned/redeemed/balance)', '3 columns = Dec 2025 totals', 'Reconstructed from transactions with obs_date cutoff'],
    ['7', 'No email_sent events', 'Entire column', 'Documented as raw count, not rate'],
    ['8', 'Null-timestamp engagement events (NaT)', '117 members (0.023%)', 'K-Means isolated into S05 (Plateau Cruiser) as artefact cluster'],
]
add_table(sl, 0.65, 3.95, 8.7, 1.15, dq, [0.2, 1.8, 1.5, 5.2], font_size=5.5)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 — Architecture & Technical Design (Light Theme)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, LIGHT_BG); add_footer(sl, 2)
add_tb(sl, 0.5, 0.25, 7, 0.5, 'Architecture & Technical Design', sz=28, bold=True, color=DARK_TEXT, font='Georgia')

# ── PIPELINE FLOW DIAGRAM ──
add_rect(sl, 0.5, 0.85, 9.0, 2.7, fill=None, line=LIGHT_ORANGE, thick_top=True)

# Raw Data box
add_rect(sl, 2.5, 0.95, 5.0, 0.3, fill=WHITE, line=ACCENT_ORANGE)
add_tb(sl, 2.6, 0.97, 4.8, 0.25, 'RAW DATA: members (500K) | transactions (17.8M) | engagement (35.5M)', sz=7, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_tb(sl, 4.8, 1.25, 0.4, 0.15, '▼', sz=10, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

# Step 01: Raw Validation
add_rect(sl, 0.7, 1.38, 8.6, 0.38, fill=WHITE, line=LIGHT_ORANGE)
add_ml(sl, 0.8, 1.4, 8.4, 0.34, [
    ('[01] Raw Validation (src/01_validate_raw.py)', True, ACCENT_ORANGE, 7),
    ('88,717 ghost IDs excluded | 92,748 duplicate events deduped | Dual-pass datetime parser | 24,738 mixed-case normalised | Leakage columns flagged', False, BODY_TEXT, 6),
])
add_tb(sl, 4.8, 1.76, 0.4, 0.12, '▼', sz=9, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

# Step 02-03: Spine + Snapshot
add_rect(sl, 0.7, 1.88, 4.15, 0.38, fill=WHITE, line=LIGHT_ORANGE)
add_ml(sl, 0.8, 1.9, 4.0, 0.34, [
    ('[02] Member Spine', True, DARK_TEXT, 7),
    ('500K real members | tenure_days clipped ≥ 0', False, BODY_TEXT, 6),
])
add_rect(sl, 5.15, 1.88, 4.15, 0.38, fill=WHITE, line=LIGHT_ORANGE)
add_ml(sl, 5.25, 1.9, 4.0, 0.34, [
    ('[03] Snapshot Builder', True, DARK_TEXT, 7),
    ('4 windows: 7d/30d/90d/180d | session cap 14,400s | points reconstructed', False, BODY_TEXT, 6),
])
add_tb(sl, 4.8, 2.26, 0.4, 0.12, '▼', sz=9, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

# Step 04-05: Feature Engineering + Validation
add_rect(sl, 0.7, 2.38, 4.15, 0.38, fill=WHITE, line=LIGHT_ORANGE)
add_ml(sl, 0.8, 2.4, 4.0, 0.34, [
    ('[04] Feature Engineering (119 features)', True, ACCENT_ORANGE, 7),
    ('OLS spend slope | 6 velocity features | ±50 clip | Cached to parquet', False, BODY_TEXT, 6),
])
add_rect(sl, 5.15, 2.38, 4.15, 0.38, fill=WHITE, line=LIGHT_ORANGE)
add_ml(sl, 5.25, 2.4, 4.0, 0.34, [
    ('[05] Feature Validation', True, ACCENT_ORANGE, 7),
    ('Range + nullity checks on all 119 features | Hard-stop on violations', False, BODY_TEXT, 6),
])

# Arrows to 3 parallel steps
add_tb(sl, 1.8, 2.76, 0.4, 0.12, '▼', sz=9, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
add_tb(sl, 4.8, 2.76, 0.4, 0.12, '▼', sz=9, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)
add_tb(sl, 7.8, 2.76, 0.4, 0.12, '▼', sz=9, color=ACCENT_ORANGE, align=PP_ALIGN.CENTER)

# Step 06, 07, 08 — Three parallel boxes
add_rect(sl, 0.7, 2.88, 2.8, 0.6, fill=WHITE, line=LIGHT_ORANGE)
add_ml(sl, 0.8, 2.9, 2.6, 0.55, [
    ('[06] Segment Assignment', True, DARK_TEXT, 6.5),
    ('FROZEN K-Means k=5', False, ACCENT_ORANGE, 6),
    ('40 features → Scaler → PCA(18)', False, BODY_TEXT, 5.5),
    ('85% variance explained', False, BODY_TEXT, 5.5),
])
add_rect(sl, 3.6, 2.88, 2.8, 0.6, fill=WHITE, line=LIGHT_ORANGE)
add_ml(sl, 3.7, 2.9, 2.6, 0.55, [
    ('[07] State Mapping', True, DARK_TEXT, 6.5),
    ('10-rule priority cascade', False, ACCENT_ORANGE, 6),
    ('numpy.select, <1 sec', False, BODY_TEXT, 5.5),
    ('supporting_evidence per member', False, BODY_TEXT, 5.5),
])
add_rect(sl, 6.5, 2.88, 2.8, 0.6, fill=WHITE, line=LIGHT_ORANGE)
add_ml(sl, 6.6, 2.9, 2.6, 0.55, [
    ('[08] Transition Prediction', True, DARK_TEXT, 6.5),
    ('FROZEN XGBoost, 49 features', False, ACCENT_ORANGE, 6),
    ('multi:softprob, 5 probs/member', False, BODY_TEXT, 5.5),
    ('Walk-forward training', False, BODY_TEXT, 5.5),
])

# ── KEY TECH CHOICES | DESIGN DECISIONS | ERROR HANDLING ──
add_rect(sl, 0.5, 3.7, 3.0, 1.5, fill=None, line=LIGHT_ORANGE)
add_tb(sl, 0.65, 3.75, 2.5, 0.18, 'KEY TECH CHOICES', sz=8, bold=True, color=ACCENT_ORANGE)
add_ml(sl, 0.65, 3.95, 2.7, 1.2, [
    ('K-Means (sklearn 1.9): Hard assignment, seed=42. Correct for continuous data where HDBSCAN produces 71–85% noise.', False, BODY_TEXT, 5.5),
    ('PCA (18 comp, 85% var): 40 correlated features make Euclidean distance unreliable. PCA used ONLY for K-Means.', False, BODY_TEXT, 5.5),
    ('XGBoost 3.3 (multi:softprob): Full 5-class probability. scale_pos_weight handles 1,700:1 imbalance.', False, BODY_TEXT, 5.5),
    ('numpy.select: State cascade <1s on 500K. Every condition is a readable business rule.', False, BODY_TEXT, 5.5),
    ('pyarrow 24.0: Feature cache ~90s vs ~22 min full rebuild.', False, BODY_TEXT, 5.5),
    ('scipy cdist: Batch pairwise distance to 5 centroids, 4× faster than loop.', False, BODY_TEXT, 5.5),
])

add_rect(sl, 3.6, 3.7, 3.0, 1.5, fill=None, line=LIGHT_ORANGE)
add_tb(sl, 3.75, 3.75, 2.5, 0.18, 'DESIGN DECISIONS', sz=8, bold=True, color=ACCENT_ORANGE)
add_ml(sl, 3.75, 3.95, 2.7, 1.2, [
    ('k=5 vs k=7: k=5: CH 77,123, F1 0.8138. k=7: 2 centroids within 0.4 PCA dist — identical.', False, BODY_TEXT, 5.5),
    ('Silhouette 0.12 — kept. Log1p+RobustScaler → sil 0.37, but F1 dropped to 0.69. Better geometry = worse predictions.', False, BODY_TEXT, 5.5),
    ('PCA before K-Means: Without PCA, spend windows quadruple-weight spending. Silent Acc merged into Prog Skeptic.', False, BODY_TEXT, 5.5),
    ('Frozen models: fit() never called at inference. Dec 2025 results reproducible to the bit.', False, BODY_TEXT, 5.5),
    ('Rule-based states: "recency=74, app_open=3" is actionable. Learned classifier is not.', False, BODY_TEXT, 5.5),
])

add_rect(sl, 6.7, 3.7, 2.8, 1.5, fill=None, line=LIGHT_ORANGE)
add_tb(sl, 6.85, 3.75, 2.5, 0.18, 'ERROR HANDLING', sz=8, bold=True, color=ACCENT_ORANGE)
add_ml(sl, 6.85, 3.95, 2.5, 1.2, [
    ('• No purchases: fillna(999) → routes to Win-Back / Lapse Risk', False, BODY_TEXT, 5.5),
    ('• Negative tenure: clip(lower=0) → day-0 member', False, BODY_TEXT, 5.5),
    ('• Slope undefined (<2 data pts): returns 0.0', False, BODY_TEXT, 5.5),
    ('• Missing cache: rebuilds prior features in-memory', False, BODY_TEXT, 5.5),
    ('• Ghost IDs: excluded at spine build, never enter features', False, BODY_TEXT, 5.5),
    ('• Mixed datetime: dual-pass parser, 0 unparseable across 17.8M rows', False, BODY_TEXT, 5.5),
    ('• NaT events: 117 members isolated into S05', False, BODY_TEXT, 5.5),
    ('• Confidence=0.5: valid output (equidistant), clipped [0,1]', False, BODY_TEXT, 5.5),
    ('• Row mismatch: assert n_rows == len(spine), hard fail', False, BODY_TEXT, 5.5),
])


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 — Results & Key Outputs (Light Theme)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, LIGHT_BG); add_footer(sl, 3)
add_tb(sl, 0.5, 0.25, 7, 0.5, 'Results & Key Outputs', sz=28, bold=True, color=DARK_TEXT, font='Georgia')

# ── LEFT: SEGMENTS DISCOVERED ──
add_tb(sl, 0.5, 0.85, 4, 0.25, 'SEGMENTS DISCOVERED (December 31, 2025)', sz=10, bold=True, color=ACCENT_ORANGE)
add_rect(sl, 0.5, 1.1, 4.8, 2.5, fill=None, line=LIGHT_ORANGE)
seg = [
    ['Segment', 'Size', '%', '30d Spend', '30d Purch', 'Key Trait'],
    ['Growth Builder', '198,035', '39.6%', '$247.10', '4.8', 'Mid-tier, multi-channel, 16,240 pts'],
    ['High-Tier Acc.', '87,688', '17.5%', '$689.30', '7.6', 'Top spenders, 61.4% PLCC'],
    ['Program Skeptic', '87,505', '17.5%', '$187.30', '4.0', 'Active buyers, tier=0, zero tier changes'],
    ['Silent Accum.', '126,655', '25.3%', '$106.00', '2.3', '5+ categories, zero digital engagement'],
    ['Plateau Cruiser', '117', '0.02%', '$307.60', '4.5', '375 app opens/30d, 10:1 browse:buy'],
]
add_table(sl, 0.6, 1.15, 4.6, 1.0, seg, [0.9, 0.5, 0.35, 0.55, 0.45, 1.85], font_size=5.5)

add_ml(sl, 0.6, 2.25, 4.6, 1.2, [
    ('Why k=5?', True, DARK_TEXT, 7),
    ('Highest Calinski-Harabasz (77,123). k=6 produced 2 unstable microclusters → F1 dropped by 0.12.', False, BODY_TEXT, 6),
    ('', False, BODY_TEXT, 3),
    ('Two segments that surprised us:', True, ACCENT_ORANGE, 7),
    ('Plateau Cruiser (117): Expected high-engagement cluster. Found data artefact — null-timestamp events bypassing 30d filter. K-Means correctly quarantined them.', False, BODY_TEXT, 6),
    ('Program Skeptic (87,505): 17.5% spends $187/mo but never engaged with any loyalty mechanic. Tier 0, zero tier changes. The program means nothing to them.', False, BODY_TEXT, 6),
])

# ── RIGHT: MODEL PERFORMANCE ──
add_tb(sl, 5.5, 0.85, 4, 0.25, 'MODEL PERFORMANCE', sz=10, bold=True, color=ACCENT_ORANGE)
add_rect(sl, 5.5, 1.1, 4.0, 2.5, fill=None, line=LIGHT_ORANGE)

add_ml(sl, 5.6, 1.15, 3.8, 2.4, [
    ('Cluster Quality (K-Means):', True, DARK_TEXT, 7.5),
    ('• Calinski-Harabasz: 77,123 (Extremely high — 5 centroids distinctly separated)', False, BODY_TEXT, 6),
    ('• Davies-Bouldin: 2.000 (Moderate boundary overlap)', False, BODY_TEXT, 6),
    ('• Silhouette: 0.120 (Low but correct — HDBSCAN proved no density gaps, 71–85% noise. We are partitioning a gradient, not finding blobs)', False, BODY_TEXT, 6),
    ('', False, BODY_TEXT, 3),
    ('Transition Prediction (XGBoost, Nov→Dec):', True, DARK_TEXT, 7.5),
    ('• Macro F1 Score: 0.8138', False, GREEN_OK, 7),
    ('• High-Tier Accelerator F1: 0.898 (highly predictable)', False, BODY_TEXT, 6),
    ('• Silent Accumulator F1: 0.666 (hardest — no digital footprint)', False, BODY_TEXT, 6),
    ('', False, BODY_TEXT, 3),
    ('State Mapping Coherence:', True, DARK_TEXT, 7.5),
    ('Priority cascade (numpy.select) cleanly resolves conflicting signals. E.g., spend_slope=3.5 → Momentum Builder overrides Plateau Cruiser via Priority 4 > Priority 9.', False, BODY_TEXT, 6),
])

# ── BOTTOM: ALGORITHM COMPARISON TABLE ──
add_rect(sl, 0.5, 3.75, 9.0, 1.45, fill=None, line=LIGHT_ORANGE, thick_top=True)
add_tb(sl, 0.65, 3.82, 4, 0.18, 'ALGORITHM COMPARISON (50K Subsample)', sz=8, bold=True, color=ACCENT_ORANGE)
algo = [
    ['Algorithm', 'Result', 'Reason Rejected'],
    ['HDBSCAN', '71–85% noise', 'No density gaps in loyalty behavioural space'],
    ['GMM', 'Silhouette 0.018, CH 32,390', 'Worst on all three metrics'],
    ['Bisecting K-Means', 'Silhouette 0.089, CH 58,236', 'Worse than K-Means on all metrics'],
    ['BIRCH', 'OOM + 64.7% in one cluster', 'Collapsed'],
    ['K-Means k=6', 'Two microclusters, −0.12 F1', 'Downstream accuracy collapsed'],
    ['K-Means k=5', 'CH 77,123, F1 0.8138', '✓ SELECTED'],
]
add_table(sl, 0.65, 4.02, 8.7, 1.0, algo, [1.2, 2.0, 5.5], font_size=5.5)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 — Business Impact & Activation (Light Theme)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, LIGHT_BG); add_footer(sl, 4)
add_tb(sl, 0.5, 0.25, 7, 0.5, 'Business Impact & Activation', sz=28, bold=True, color=DARK_TEXT, font='Georgia')

# ── LEFT: ACTIVATION STRATEGY ──
add_tb(sl, 0.5, 0.85, 4, 0.25, 'ACTIVATION STRATEGY (Segment × State)', sz=10, bold=True, color=ACCENT_ORANGE)
add_rect(sl, 0.5, 1.1, 4.3, 4.1, fill=None, line=LIGHT_ORANGE)

# Example 1
add_ml(sl, 0.6, 1.15, 4.1, 1.0, [
    ('❶ High-Tier Accelerator in Lapse Risk', True, RED_ALERT, 7.5),
    ('(Historically $689/mo, 0 purchases in last 30d, spend slope declining)', False, BODY_TEXT, 6),
    ('Channel: Personal email + PLCC Account Manager (61.4% hold the card)', False, DARK_TEXT, 6),
    ('Message: "Exclusive preview — reserved for Platinum members."', False, DARK_TEXT, 6),
    ('Offer: Early access to new product line. No discounts.', False, DARK_TEXT, 6),
    ('Timing: Within 7 days of 30-day lapse threshold.', False, DARK_TEXT, 6),
    ('Why: They respond to VIP treatment, not $5 coupons. (Tone: VIP/Exclusive)', False, ACCENT_ORANGE, 6),
])

# Example 2
add_ml(sl, 0.6, 2.2, 4.1, 1.0, [
    ('❷ Program Skeptic in Momentum Builder', True, DARK_TEXT, 7.5),
    ('(Ignores program, but spend_slope > 2.0 with recent purchases)', False, BODY_TEXT, 6),
    ('Channel: App Push (capitalise on current momentum)', False, DARK_TEXT, 6),
    ('Message: "You\'re 1 purchase away from Gold status."', False, DARK_TEXT, 6),
    ('Offer: Tier accelerator bonus. Timing: Immediate.', False, DARK_TEXT, 6),
    ('Why: Skeptics need value proof. Show them they\'re already winning. (Tone: Low-commitment / Value proof)', False, ACCENT_ORANGE, 6),
])

# Example 3
add_ml(sl, 0.6, 3.25, 4.1, 0.9, [
    ('❸ Silent Accumulator in Win-Back Target', True, GREEN_OK, 7.5),
    ('($106/mo across 5 categories, zero digital footprint, lapsed 60+ days)', False, BODY_TEXT, 6),
    ('Channel: Email + SMS (app_open_30d is literally 0, push is useless)', False, DARK_TEXT, 6),
    ('Message: "We miss you." Offer: Reactivation bonus points on next in-store swipe.', False, DARK_TEXT, 6),
    ('Timing: Immediate upon crossing 60 days.', False, DARK_TEXT, 6),
    ('Why: RFM sends push they\'ll never see. TBIE routes to SMS/Email. (Tone: Simple / Transactional)', False, ACCENT_ORANGE, 6),
])

# ── RIGHT TOP: SO WHAT — THE BUSINESS CASE ──
add_tb(sl, 5.2, 0.85, 4, 0.25, 'SO WHAT — THE BUSINESS CASE', sz=10, bold=True, color=ACCENT_ORANGE)
add_rect(sl, 5.2, 1.1, 4.3, 2.0, fill=None, line=LIGHT_ORANGE)
add_ml(sl, 5.3, 1.15, 4.1, 1.85, [
    ('Why a loyalty manager should care:', True, DARK_TEXT, 7),
    ('Traditional RFM collapses trajectory. Two members spending $200 this month look identical to RFM.', False, BODY_TEXT, 6),
    ('TBIE knows one is a Growth Builder accelerating (slope=3.5) and one is a Program Skeptic decelerating (slope=-1.2). RFM sends both a generic newsletter. TBIE sends the Builder a tier-upgrade push and the Skeptic a win-back discount.', False, BODY_TEXT, 6),
    ('', False, BODY_TEXT, 3),
    ('Quantifiable Impact:', True, GREEN_OK, 7),
    ('The two-layer model is an early warning system. It detects a High-Tier Accelerator entering Lapse Risk 30 days before their 180d segment average drops.', False, BODY_TEXT, 6),
    ('Retaining just 10% of 87,688 High-Tier Accelerators at Lapse Risk preserves $6.04M in monthly revenue (8,768 × $689/mo) that would otherwise quietly churn.', False, GREEN_OK, 6.5),
])

# ── RIGHT BOTTOM: PLCC CARDHOLDER INSIGHT ──
add_rect(sl, 5.2, 3.25, 4.3, 1.95, fill=None, line=LIGHT_ORANGE, thick_top=True)
add_tb(sl, 5.35, 3.35, 4.0, 0.18, 'PLCC CARDHOLDER INSIGHT', sz=8, bold=True, color=ACCENT_ORANGE)
add_ml(sl, 5.35, 3.55, 4.0, 0.3, [
    ('Card status was NOT used to build clusters, yet segments naturally separated by card ownership based entirely on behaviour:', False, BODY_TEXT, 5.5),
])
card = [
    ['Segment', 'PLCC %', 'Behavioural Reality', 'Activation Implication'],
    ['High-Tier Acc.', '61.4%', 'Core value engine — high spend/pts/card', 'Messaging must reference card-specific benefits'],
    ['Growth Builder', '34.9%', 'The conversion battleground', 'Non-cardholders = primary PLCC acquisition targets'],
    ['Prog. Skeptic', '25.7%', 'They buy the product, not the program', 'Card offers tied to product discounts, not points'],
    ['Silent Accum.', '13.7%', 'Cash/debit, in-store, ignore digital', 'Don\'t waste digital spend. POS receipt offers only'],
]
add_table(sl, 5.35, 3.9, 4.0, 1.2, card, [0.8, 0.4, 1.1, 1.7], font_size=5.5)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 (Dark Theme) — Live Demo Highlights
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, TITLE_BG); add_footer(sl, 5)

add_tb(sl, 0.5, 0.25, 7, 0.5, 'Live Demo Highlights', sz=28, bold=True, color=WHITE, font='Georgia')
add_tb(sl, 0.5, 0.75, 7, 0.3, "Show, don't tell. Walk the judges through a live run or key screenshots.", sz=12, color=LIGHT_GRAY)

BOX_FILL = RGBColor(45, 55, 70)
CODE_BG = RGBColor(20, 25, 30)

# Box 1: Pipeline Execution
add_rect(sl, 0.5, 1.15, 9.0, 1.1, fill=BOX_FILL, line=None)
add_tb(sl, 0.6, 1.25, 4, 0.2, '1. PIPELINE EXECUTION', sz=11, bold=True, color=ACCENT_ORANGE)

add_rect(sl, 0.65, 1.5, 6.0, 0.25, fill=CODE_BG, line=None)
add_tb(sl, 0.7, 1.52, 6.0, 0.2, '> python pipeline.py --data_dir ./data/train/ --observation_date 2025-12-31 --output_dir ./outputs/', sz=7, color=GREEN_OK, font='Courier New')

add_ml(sl, 0.65, 1.8, 6.0, 0.4, [
    ('✓ Single command → all 5 output files', False, WHITE, 8),
    ('✓ Auto-installs deps | Full: ~340s | Byte-identical (seed=42)', False, WHITE, 8)
])

add_tb(sl, 6.8, 1.25, 2, 0.2, 'OUTPUT FILES:', sz=8, bold=True, color=ACCENT_TEAL)
add_ml(sl, 6.8, 1.45, 2.5, 0.7, [
    ('segment_assignments.csv (500K)', False, WHITE, 7),
    ('state_assignments.csv     (500K)', False, WHITE, 7),
    ('transition_predictions.csv(500K)', False, WHITE, 7),
    ('segment_profiles.json     (5)', False, WHITE, 7),
    ('feature_descriptions.json (119)', False, WHITE, 7)
], font='Courier New')

# Box 2: Temporal Flexibility
add_rect(sl, 0.5, 2.4, 9.0, 1.0, fill=BOX_FILL, line=None)
add_tb(sl, 0.6, 2.5, 4, 0.2, '2. TEMPORAL FLEXIBILITY', sz=11, bold=True, color=ACCENT_ORANGE)

add_rect(sl, 0.65, 2.75, 4.0, 0.25, fill=CODE_BG, line=None)
add_tb(sl, 0.7, 2.77, 4.0, 0.2, '> python pipeline.py --observation_date 2025-04-29', sz=7, color=GREEN_OK, font='Courier New')

add_rect(sl, 4.8, 2.75, 4.0, 0.25, fill=CODE_BG, line=None)
add_tb(sl, 4.85, 2.77, 4.0, 0.2, '> python pipeline.py --observation_date 2025-08-31', sz=7, color=GREEN_OK, font='Courier New')

add_ml(sl, 0.65, 3.05, 8.0, 0.3, [
    ('✓ All dates dynamic — NOTHING hardcoded | Frozen centroids (transform-only) | Future data excluded', False, WHITE, 8),
    ('✓ Distribution shifts naturally: more "New & Uncertain" in early months (shorter tenure)', False, WHITE, 8)
])

# Box 3: Output Walkthrough
add_rect(sl, 0.5, 3.55, 9.0, 1.6, fill=BOX_FILL, line=None)
add_tb(sl, 0.6, 3.65, 8, 0.2, '3. OUTPUT WALKTHROUGH — Trace Member MBR_0000004 End-to-End', sz=11, bold=True, color=ACCENT_ORANGE)

trace = [
    ['Output File','Field','Value','Interpretation'],
    ['segment_assignments','segment_name','Growth Builder (S01)','Active buyer, upward trajectory'],
    ['segment_assignments','confidence','0.5849','Moderate — near S02 boundary'],
    ['state_assignments','state_name','Momentum Builder','spend_slope:29.2, 5 purch/30d'],
    ['transition_pred','prob_S02','0.4442','Nearly upgrading to High-Tier'],
]
add_table(sl, 0.65, 3.9, 8.7, 0.85, trace, [2.0, 1.5, 2.0, 3.2], font_size=7, theme='dark')

add_rect(sl, 0.65, 4.85, 8.7, 0.2, fill=RGBColor(20, 50, 20), line=GREEN_OK)
add_tb(sl, 0.7, 4.87, 8.5, 0.15, '→ ACTION: Send tier upgrade nudge immediately. Spend slope 29.2 (extreme momentum), Tier 3, 44.4% chance of upgrading to HTA next month.', sz=7, bold=True, color=GREEN_OK)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 — Learnings, Limitations & Next Steps (Light Theme)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl, LIGHT_BG); add_footer(sl, 6)
add_tb(sl, 0.5, 0.25, 7, 0.5, 'Learnings, Limitations & Next Steps', sz=28, bold=True, color=DARK_TEXT, font='Georgia')

# ── LEFT: WHAT WE LEARNED ──
add_tb(sl, 0.5, 0.85, 4, 0.25, 'WHAT WE LEARNED', sz=11, bold=True, color=ACCENT_ORANGE)
add_rect(sl, 0.5, 1.1, 4.3, 2.5, fill=None, line=LIGHT_ORANGE)
add_ml(sl, 0.65, 1.15, 4.0, 2.4, [
    ('What surprised us:', True, DARK_TEXT, 7.5),
    ('The Loyalty Gradient: Members don\'t form neat groups. HDBSCAN proved this by labelling 71–85% as noise. Segmentation is partitioning a continuum into actionable zones.', False, BODY_TEXT, 6),
    ('The Plateau Cruiser Artefact: 117 members with 375 app opens/month and 10:1 browse-to-buy. Traced to NaT event rows bypassing 30d filters. K-Means correctly isolated the dirty data.', False, BODY_TEXT, 6),
    ('', False, BODY_TEXT, 3),
    ('Biggest impact decision:', True, ACCENT_ORANGE, 7.5),
    ('Dual-Layer Architecture (Segment + State). Trying to do both with one model failed. K-Means handles slow-moving structural identity. 10-rule cascade handles fast-moving posture. XGBoost predicts the intersection. Separating concerns made the system usable for campaign managers.', False, BODY_TEXT, 6),
    ('', False, BODY_TEXT, 3),
    ('What we\'d tell another team:', True, DARK_TEXT, 7.5),
    ('Fix leakage first. members.parquet contained lifetime points reflecting Dec totals. If we hadn\'t reconstructed from raw transactions, the model would have leaked the future, scored perfectly in testing, and failed in production.', False, BODY_TEXT, 6),
])

# ── RIGHT TOP: HONEST LIMITATIONS ──
add_tb(sl, 5.2, 0.85, 4, 0.25, 'HONEST LIMITATIONS', sz=11, bold=True, color=ACCENT_ORANGE)
add_rect(sl, 5.2, 1.1, 4.3, 2.0, fill=None, line=LIGHT_ORANGE)
add_ml(sl, 5.35, 1.15, 4.0, 1.9, [
    ('Where the model struggles:', True, RED_ALERT, 7.5),
    ('• Program Skeptic is a catch-all (87,505 / 17.5%). Likely contains multiple sub-postures our 10-rule cascade collapses together.', False, BODY_TEXT, 6),
    ('• Fixed Thresholds: State engine uses hardcoded rules (e.g., recency > 60). Cannot adapt dynamically to seasonal shifts without manual recalibration.', False, BODY_TEXT, 6),
    ('• Extrapolation beyond Month 12: XGBoost trained on pairs up to Nov→Dec. Jan 2026 predictions are one month outside training distribution.', False, BODY_TEXT, 6),
    ('', False, BODY_TEXT, 3),
    ('Unresolved assumptions:', True, RED_ALERT, 7.5),
    ('• Frozen Centroids: Assumes Dec 2025 centroids remain valid. Major program restructure → severe assignment drift.', False, BODY_TEXT, 6),
    ('• Causal Inference: We predict where a member is going, but cannot quantify exact counterfactual lift of interventions.', False, BODY_TEXT, 6),
])

# ── LEFT BOTTOM: INDIVIDUAL CONTRIBUTIONS ──
add_rect(sl, 0.5, 3.75, 4.3, 1.5, fill=None, line=LIGHT_ORANGE, thick_top=True)
add_tb(sl, 0.65, 3.82, 4.0, 0.18, 'INDIVIDUAL CONTRIBUTIONS', sz=8, bold=True, color=ACCENT_ORANGE)
con_tb = [
    ['Member', 'Phases', 'Contribution'],
    ['Samarth Vinod Hosalli\n(PES1UG23AM261)', 'Phase 6 & 7', 'K-Means clustering pipeline (k=5 selection, PCA reduction, algorithm evaluation), 10-state priority cascade rule engine, methodology & data quality reports'],
    ['Siddarth Reddy\n(PES1UG23AM300)', 'Phase 1–5 & 8', 'Data validation, leakage prevention, 119-feature engineering pipeline, XGBoost walk-forward transition model, pipeline.py single-command architecture'],
]
add_table(sl, 0.65, 4.05, 4.0, 1.1, con_tb, [1.2, 0.6, 2.2], font_size=5.5)

# ── RIGHT BOTTOM: IF WE HAD MORE TIME ──
add_rect(sl, 5.2, 3.25, 4.3, 1.9, fill=None, line=LIGHT_ORANGE, thick_top=True)
add_tb(sl, 5.35, 3.32, 4.0, 0.18, 'IF WE HAD MORE TIME', sz=8, bold=True, color=ACCENT_ORANGE)
add_ml(sl, 5.35, 3.52, 4.0, 1.5, [
    ('Hidden Markov Model (HMM) for States:', True, DARK_TEXT, 7),
    ('Replace hardcoded 10-rule cascade with HMM trained on 11-month member sequences. Probabilistic state assignments that adapt to seasonal shifts automatically.', False, BODY_TEXT, 6),
    ('', False, BODY_TEXT, 3),
    ('Centroid Drift Monitoring:', True, DARK_TEXT, 7),
    ('Automated tracking of Calinski-Harabasz month-over-month. If compactness degrades below threshold, trigger alert for K-Means retraining.', False, BODY_TEXT, 6),
    ('', False, BODY_TEXT, 3),
    ('Causal Impact Estimation:', True, DARK_TEXT, 7),
    ('Integrate historical A/B test data to compute predicted counterfactual lift. Move from "predicting transitions" to "optimizing interventions."', False, BODY_TEXT, 6),
])


# ══════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════
out = r'c:\Users\sidda\Downloads\TBIE_CODE\TBIE_Solution_Presentation.pptx'
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
print("Done!")
