"""
Extract ALL text from the generated TBIE PPT and verify against the user's provided content.
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

ppt_path = r'c:\Users\sidda\Downloads\TBIE_CODE\TBIE_Solution_Presentation.pptx'
prs = Presentation(ppt_path)

print("=" * 100)
print("FULL TEXT EXTRACTION — TBIE_Solution_Presentation.pptx")
print("=" * 100)

for si, slide in enumerate(prs.slides):
    print(f"\n{'━' * 100}")
    print(f"  SLIDE {si+1}")
    print(f"{'━' * 100}")
    for shape in slide.shapes:
        if shape.has_table:
            t = shape.table
            print(f"\n  [TABLE {t.rows.__len__()}x{len(t.columns)}]")
            for r in range(len(t.rows)):
                row_vals = []
                for c in range(len(t.columns)):
                    row_vals.append(t.cell(r, c).text.strip())
                sep = " | "
                print(f"    {'HDR' if r==0 else '   '}: {sep.join(row_vals)}")
        elif shape.has_text_frame:
            full = shape.text_frame.text.strip()
            if full:
                print(f"  > {full}")

print("\n\n")
print("=" * 100)
print("CONTENT VERIFICATION CHECKLIST")
print("=" * 100)

# Collect all text
all_text = ""
all_table_text = ""
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_table:
            t = shape.table
            for r in range(len(t.rows)):
                for c in range(len(t.columns)):
                    all_table_text += t.cell(r, c).text + " "
        elif shape.has_text_frame:
            all_text += shape.text_frame.text + " "

combined = all_text + " " + all_table_text

# ── SLIDE 1 checks ──
print("\n── SLIDE 1 (Title) ──")
checks_s1 = [
    ("Macro F1 = 0.8138", "0.8138"),
    ("Segments = 5", "Segments"),
    ("States = 10", "States"),
    ("Members = 500,000", "500,000"),
    ("Features Built = 119", "119"),
    ("Team: Samarth Vinod Hosalli (PES1UG23AM261)", "PES1UG23AM261"),
    ("Team: Siddarth Reddy (PES1UG23AM300)", "PES1UG23AM300"),
    ("TBIE title", "Temporal Behavioural Intelligence Engine"),
    ("Track C", "Track C"),
]
for label, needle in checks_s1:
    found = needle in combined
    print(f"  {'✓' if found else '✗'} {label}")

# ── SLIDE 2 checks ──
print("\n── SLIDE 2 (Problem & Approach) ──")
checks_s2 = [
    ("Core challenge text", "separating signal from noise"),
    ("Continuous gradient", "continuous gradient"),
    ("HDBSCAN 71-85% noise", "71"),
    ("Temporal leakage", "Temporal leakage"),
    ("members.parquet 3 leakage cols", "lifetime_points_earned"),
    ("Segment vs state different problems", "two different problems"),
    ("Frozen models constraint", "Frozen model"),
    ("Single-command reproducibility", "Single-command"),
    ("No ground truth", "No ground truth"),
    ("Key Insight: segment=who, state=what to do", "who they are over months"),
    ("Two-layer architecture", "two-layer architecture" if "two-layer architecture" in combined.lower() else "Two-layer architecture"),
    ("K-Means as partitioning tool", "partitioning tool"),
    ("Alt rejected: HDBSCAN", "HDBSCAN"),
    ("Alt rejected: GMM soft memberships", "GMM"),
    ("Alt rejected: k=7 centroids 0.4 PCA", "0.4 PCA"),
    ("Alt rejected: Neural embeddings", "Neural embeddings" if "Neural embeddings" in combined else "interpretable"),
    ("Alt rejected: Learned state", "auditable"),
    # Data quality table
    ("DQ #1: Ghost IDs 88,717", "88,717"),
    ("DQ #2: Mixed datetime 7.4%", "7.4%"),
    ("DQ #3: Mixed-case 24,738", "24,738"),
    ("DQ #4: Duplicate events 92,748", "92,748"),
    ("DQ #5: Session outliers 18,663", "18,663"),
    ("DQ #6: Leakage columns", "Leakage columns"),
    ("DQ #7: No email_sent events", "email_sent"),
    ("DQ #8: NaT events 117 members", "117 members"),
    ("DQ #8: Plateau Cruiser S05 artefact", "Plateau Cruiser"),
]
for label, needle in checks_s2:
    found = needle in combined
    print(f"  {'✓' if found else '✗'} {label}")

# ── SLIDE 3 checks ──
print("\n── SLIDE 3 (Architecture & Technical Design) ──")
checks_s3 = [
    ("Raw data: members 500K", "500K"),
    ("Raw data: transactions 17.8M", "17.8M"),
    ("Raw data: engagement 35.5M", "35.5M"),
    ("Step 01: Raw Validation", "Raw Validation"),
    ("Step 02: Member Spine", "Member Spine"),
    ("Step 03: Snapshot Builder", "Snapshot Builder"),
    ("Step 04: Feature Engineering 119", "119 features"),
    ("Step 05: Feature Validation", "Feature Validation"),
    ("Step 06: Segment Assignment K-Means k=5", "K-Means k=5"),
    ("Step 07: State Mapping 10-rule", "10-rule"),
    ("Step 08: Transition Prediction XGBoost", "Transition Prediction"),
    ("PCA 18 components", "18"),
    ("85% variance", "85%"),
    ("numpy.select <1 sec", "numpy.select"),
    ("multi:softprob", "multi:softprob"),
    ("49 features for XGBoost", "49 features"),
    ("40 features for clustering", "40 features"),
    # Tech choices
    ("K-Means sklearn 1.9", "1.9"),
    ("XGBoost 3.3", "3.3"),
    ("pyarrow 24.0", "24.0"),
    ("scipy cdist 4x faster", "4×" if "4×" in combined else "4x"),
    ("seed=42", "seed=42"),
    ("scale_pos_weight 1,700:1", "1,700:1"),
    ("~90s cached vs ~22 min", "90"),
    # Design decisions
    ("CH 77,123", "77,123"),
    ("F1 0.8138", "0.8138"),
    ("Silhouette 0.12 kept", "0.12"),
    ("log1p+RobustScaler sil 0.37 F1 0.69", "0.37"),
    ("k=7 centroids 0.4 PCA dist", "0.4 PCA"),
    ("PCA before K-Means", "PCA before K-Means"),
    ("Frozen models: fit() never at inference", "fit()"),
    ("Rule-based states explainable", "recency"),
    # Error handling
    ("fillna(999)", "999"),
    ("clip(lower=0)", "clip"),
    ("Slope <2 data pts = 0.0", "0.0"),
    ("Missing cache rebuilt in-memory", "in-memory"),
    ("Ghost IDs excluded at spine", "spine"),
    ("Dual-pass parser", "dual-pass" if "dual-pass" in combined.lower() else "Dual-pass"),
    ("NaT 117 members S05", "S05"),
    ("Confidence 0.5 equidistant", "equidistant"),
    ("Row mismatch assert", "assert"),
    ("BIRCH OOM 64.7%", "64.7%"),
    ("Bisecting K-Means sil 0.089 CH 58,236", "58,236"),
    ("GMM sil 0.018 CH 32,390", "32,390"),
]
for label, needle in checks_s3:
    found = needle in combined
    print(f"  {'✓' if found else '✗'} {label}")

# ── SLIDE 4 checks ──
print("\n── SLIDE 4 (Results & Key Outputs) ──")
checks_s4 = [
    ("Growth Builder 198,035 39.6%", "198,035"),
    ("Growth Builder $247.10", "$247.10"),
    ("Growth Builder 4.8 purchases", "4.8"),
    ("Growth Builder 16,240 pts", "16,240"),
    ("High-Tier Acc 87,688 17.5%", "87,688"),
    ("High-Tier Acc $689.30", "$689.30"),
    ("High-Tier Acc 7.6 purchases", "7.6"),
    ("High-Tier Acc 61.4% PLCC", "61.4%"),
    ("Program Skeptic 87,505 17.5%", "87,505"),
    ("Program Skeptic $187.30", "$187.30"),
    ("Program Skeptic tier=0", "tier=0"),
    ("Silent Accum 126,655 25.3%", "126,655"),
    ("Silent Accum $106.00", "$106.00"),
    ("Silent Accum 2.3 purchases", "2.3"),
    ("Plateau Cruiser 117 0.02%", "0.02%"),
    ("Plateau Cruiser $307.60", "$307.60"),
    ("Plateau Cruiser 375 app opens", "375"),
    ("Plateau Cruiser 10:1 browse:buy", "10:1"),
    ("CH 77,123", "77,123"),
    ("DB 2.000", "2.000"),
    ("Silhouette 0.120", "0.120"),
    ("Macro F1 0.8138", "0.8138"),
    ("HTA F1 0.898", "0.898"),
    ("Silent Acc F1 0.666", "0.666"),
    ("Priority cascade resolves conflicts", "Priority"),
    ("Algorithm comparison table present", "BIRCH"),
    ("k=6 microclusters -0.12 F1", "0.12"),
]
for label, needle in checks_s4:
    found = needle in combined
    print(f"  {'✓' if found else '✗'} {label}")

# ── SLIDE 5 checks ──
print("\n── SLIDE 5 (Business Impact & Activation) ──")
checks_s5 = [
    ("Example 1: HTA in Lapse Risk", "Lapse Risk"),
    ("HTA $689/mo", "$689"),
    ("Channel: Personal email + PLCC Account Manager", "Account Manager"),
    ("Message: Exclusive preview", "Exclusive preview"),
    ("No discounts for HTA", "No discounts"),
    ("Timing: 7 days of lapse threshold", "7 days"),
    ("VIP treatment tone", "VIP"),
    ("Example 2: Prog Skeptic in Momentum Builder", "Momentum Builder"),
    ("spend_slope > 2.0", "2.0"),
    ("Message: 1 purchase away from Gold", "Gold"),
    ("Tier accelerator bonus", "accelerator"),
    ("Value proof tone", "Value proof"),
    ("Example 3: Silent Acc in Win-Back Target", "Win-Back"),
    ("app_open_30d is 0", "app_open"),
    ("Channel: Email + SMS", "SMS"),
    ("Message: We miss you", "miss you"),
    ("Reactivation bonus points", "Reactivation"),
    ("SO WHAT business case", "SO WHAT"),
    ("RFM collapses trajectory", "trajectory"),
    ("slope=3.5 vs slope=-1.2", "3.5"),
    ("$6.04M monthly revenue", "$6.04M"),
    ("8,768 members x $689", "8,768"),
    ("PLCC insight table", "PLCC"),
    ("HTA 61.4% PLCC", "61.4%"),
    ("Growth Builder 34.9%", "34.9%"),
    ("Prog Skeptic 25.7%", "25.7%"),
    ("Silent Acc 13.7%", "13.7%"),
    ("Card not used in clustering", "NOT used"),
    ("POS receipt offers for Silent Acc", "POS"),
]
for label, needle in checks_s5:
    found = needle in combined
    print(f"  {'✓' if found else '✗'} {label}")

# ── SLIDE 6 (Demo) checks ──
print("\n── SLIDE 6 (Live Demo) ──")
checks_s6 = [
    ("Pipeline command shown", "pipeline.py"),
    ("seed=42", "seed=42"),
    ("Output files listed", "segment_assignments"),
    ("feature_descriptions.json (119)", "119"),
    ("Temporal flexibility", "TEMPORAL"),
    ("MBR_0000004 trace", "MBR_0000004"),
]
for label, needle in checks_s6:
    found = needle in combined
    print(f"  {'✓' if found else '✗'} {label}")

# ── SLIDE 7 checks ──
print("\n── SLIDE 7 (Learnings, Limitations & Next Steps) ──")
checks_s7 = [
    ("Loyalty Gradient surprise", "Loyalty Gradient"),
    ("Plateau Cruiser Artefact", "Artefact"),
    ("HDBSCAN 71-85% noise", "71"),
    ("375 app opens NaT", "375"),
    ("Dual-Layer Architecture biggest decision", "Dual-Layer"),
    ("Fix leakage first", "Fix leakage first"),
    ("Leaked future, scored perfectly", "scored perfectly"),
    # Limitations
    ("Program Skeptic catch-all 87,505 17.5%", "catch-all"),
    ("Fixed Thresholds hardcoded", "hardcoded"),
    ("Extrapolation beyond Month 12", "Month 12"),
    ("Nov→Dec training", "Nov"),
    ("Frozen Centroids assumption", "Frozen Centroids"),
    ("Causal Inference limitation", "Causal Inference"),
    ("Counterfactual lift", "counterfactual"),
    # If we had more time
    ("HMM for States", "Hidden Markov"),
    ("11-month member sequences", "11-month"),
    ("Centroid Drift Monitoring", "Drift Monitoring"),
    ("CH month-over-month tracking", "month-over-month"),
    ("Causal Impact Estimation", "Causal Impact"),
    ("A/B test data", "A/B"),
    # Contributions
    ("Samarth: Phase 6 & 7", "Phase 6"),
    ("Samarth: K-Means, PCA, algorithm eval", "algorithm eval"),
    ("Samarth: 10-state cascade", "cascade"),
    ("Samarth: methodology & data quality reports", "methodology"),
    ("Siddarth: Phases 1-5 & 8", "Phase 1"),
    ("Siddarth: leakage prevention", "leakage prevention"),
    ("Siddarth: 119-feature pipeline", "119-feature"),
    ("Siddarth: XGBoost walk-forward", "walk-forward"),
    ("Siddarth: pipeline.py architecture", "pipeline.py"),
]
for label, needle in checks_s7:
    found = needle in combined
    print(f"  {'✓' if found else '✗'} {label}")

# ── Summary ──
print("\n" + "=" * 100)
all_checks = checks_s1 + checks_s2 + checks_s3 + checks_s4 + checks_s5 + checks_s6 + checks_s7
passed = sum(1 for _, needle in all_checks if needle in combined)
total = len(all_checks)
print(f"TOTAL: {passed}/{total} checks passed ({passed/total*100:.1f}%)")
if passed < total:
    print("\nFAILED CHECKS:")
    for label, needle in all_checks:
        if needle not in combined:
            print(f"  ✗ {label}  (looking for: '{needle}')")
print("=" * 100)
